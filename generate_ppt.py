#!/usr/bin/env python3
"""Opclaw 产品介绍 PPT 生成脚本"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 常量 ──
W, H = Inches(13.333), Inches(7.5)  # 16:9
SCREENSHOT_DIR = r"C:\Users\Administrator\.gemini\antigravity\brain\d3127ad1-9752-47f8-8e20-6178ec2058f9"
OUTPUT = r"e:\Code\AI\Start\Web\Superui\Opclaw_产品介绍.pptx"

# 配色
C_BG_DARK   = RGBColor(0x0F, 0x0A, 0x1E)
C_BG_LIGHT  = RGBColor(0x16, 0x10, 0x2A)
C_PRIMARY   = RGBColor(0x00, 0xD4, 0xFF)
C_ACCENT    = RGBColor(0xA8, 0x55, 0xF7)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_TEXT_SEC   = RGBColor(0xCB, 0xD5, 0xE1)
C_TEXT_DIM   = RGBColor(0x94, 0xA3, 0xB8)
C_ROSE      = RGBColor(0xF4, 0x3F, 0x5E)
C_EMERALD   = RGBColor(0x10, 0xB9, 0x81)
C_AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
C_BLUE      = RGBColor(0x3B, 0x82, 0xF6)
C_PINK      = RGBColor(0xEC, 0x48, 0x99)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

# ── 工具函数 ──
def dark_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG_DARK

def add_rect(slide, left, top, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, w, h, text, size=18, color=C_WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_gradient_rect(slide, left, top, w, h, color1, color2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0
    shape.line.fill.background()
    return shape

def add_screenshot(slide, filename, left, top, w, h):
    path = os.path.join(SCREENSHOT_DIR, filename)
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, left, top, w, h)
        return pic
    return None

def section_title_bar(slide, text):
    add_gradient_rect(slide, Inches(0), Inches(0), W, Inches(0.08), C_PRIMARY, C_ACCENT)
    add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.6), text, size=32, color=C_PRIMARY, bold=True)

def add_bullet_text(slide, left, top, w, h, items, size=14, color=C_TEXT_SEC):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(6)
    return txBox

# ══════════════════════════════════════════════
# SLIDE 1: 封面
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
# 装饰渐变圆
add_gradient_rect(slide, Inches(-2), Inches(-2), Inches(8), Inches(8), C_PRIMARY, C_ACCENT)
add_rect(slide, Inches(0), Inches(0), W, H, C_BG_DARK)
# 顶部装饰线
add_gradient_rect(slide, Inches(0), Inches(0), W, Inches(0.06), C_PRIMARY, C_ACCENT)
# Logo 区域
add_text(slide, Inches(0), Inches(1.2), W, Inches(0.8), "🌈", size=60, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(2.2), W, Inches(0.8), "Opclaw", size=54, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.1), W, Inches(0.5), "全能数字资产与 AI 分身助手", size=24, color=C_PRIMARY, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3.8), Inches(9.3), Inches(0.8),
         "面向 OPC 超级个体，融合个人主页 · 学习空间 · 工作助手 · 生活记录 · AI 数字人分身",
         size=16, color=C_TEXT_SEC, align=PP_ALIGN.CENTER)
# 底部标签
tags = ["React 19", "TypeScript", "Three.js", "RAG Engine", "5 套主题", "Supabase"]
tag_start = Inches(2.5)
for i, tag in enumerate(tags):
    x = tag_start + Inches(i * 1.45)
    add_rect(slide, x, Inches(5.0), Inches(1.3), Inches(0.4), RGBColor(0x25, 0x1D, 0x42))
    add_text(slide, x, Inches(5.03), Inches(1.3), Inches(0.35), tag, size=11, color=C_PRIMARY, align=PP_ALIGN.CENTER)
# 底部
add_text(slide, Inches(0), Inches(6.5), W, Inches(0.4), "一个人的数字宇宙，从这里开始！", size=14, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(6.9), W, Inches(0.3), "Made with ❤️ by 晓叶  |  2026", size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 2: 目录
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "📋 目录 CONTENTS")
chapters = [
    ("01", "产品概述", "解决什么问题 · 目标用户 · 核心价值", C_PRIMARY),
    ("02", "核心功能模块", "9 大功能模块详细介绍", C_ACCENT),
    ("03", "技术架构亮点", "前沿技术栈 · AI 能力 · 创新点", C_EMERALD),
    ("04", "用户体验设计", "5 套主题 · 响应式 · 交互体验", C_AMBER),
    ("05", "市场分析", "竞品分析 · 商业机会 · 目标市场", C_ROSE),
    ("06", "总结与展望", "项目优势 · 发展愿景 · 路线图", C_BLUE),
]
for i, (num, title, desc, color) in enumerate(chapters):
    y = Inches(1.3 + i * 0.95)
    add_gradient_rect(slide, Inches(1.5), y, Inches(10.3), Inches(0.8), RGBColor(0x1A, 0x14, 0x30), RGBColor(0x22, 0x1A, 0x3D))
    add_text(slide, Inches(1.8), y + Inches(0.05), Inches(0.8), Inches(0.7), num, size=28, color=color, bold=True)
    add_text(slide, Inches(2.8), y + Inches(0.08), Inches(3), Inches(0.4), title, size=20, color=C_WHITE, bold=True)
    add_text(slide, Inches(2.8), y + Inches(0.42), Inches(6), Inches(0.3), desc, size=12, color=C_TEXT_DIM)
    # 右侧装饰线
    add_rect(slide, Inches(11.5), y + Inches(0.15), Inches(0.06), Inches(0.5), color)

# ══════════════════════════════════════════════
# SLIDE 3: 产品概述 - 痛点
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "01  产品概述 — 用户痛点与解决方案")
add_text(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
         "随着 AI 技术爆发与个体经济崛起，OPC 超级个体面临五大核心矛盾：", size=15, color=C_TEXT_SEC)

pains = [
    ("😖 个人 IP 展示分散", "简历在 Word、作品在 GitHub、社交在各平台", "一站式个人主页 + 在线简历 + 作品集 + 社交矩阵", "减少 80% 展示链接", C_ROSE),
    ("📝 知识管理碎片化", "笔记散落 Notion / 语雀 / 本地文件", "知识库三栏布局 + 富文本编辑 + AI 问答", "检索效率提升 3x", C_AMBER),
    ("🤖 AI 能力未个性化", "通用 AI 助手不懂我的知识体系", "自研 RAG 引擎 + 声音/形象克隆", "回答相关度提升 5x", C_ACCENT),
    ("💔 生活记忆无归处", "恋爱纪念靠截图、旅行照片在相册", "恋爱时间轴 + 地图轨迹 + 时光相册", "记忆完整度提升 90%", C_PINK),
    ("🔧 工作工具割裂", "新媒体 / 电商 / 工具分别在不同 SaaS", "新媒体 + 电商 + 百宝箱一体化", "切换成本降低 70%", C_EMERALD),
]
for i, (pain, current, solution, value, color) in enumerate(pains):
    y = Inches(1.6 + i * 1.1)
    add_rect(slide, Inches(0.6), y, Inches(12), Inches(0.95), RGBColor(0x1A, 0x14, 0x30))
    add_text(slide, Inches(0.8), y + Inches(0.05), Inches(2.5), Inches(0.35), pain, size=14, color=color, bold=True)
    add_text(slide, Inches(0.8), y + Inches(0.4), Inches(2.5), Inches(0.5), "现状: " + current, size=10, color=C_TEXT_DIM)
    add_text(slide, Inches(3.5), y + Inches(0.15), Inches(5.5), Inches(0.6), "→  " + solution, size=12, color=C_WHITE)
    add_rect(slide, Inches(9.5), y + Inches(0.2), Inches(2.8), Inches(0.5), RGBColor(0x25, 0x1D, 0x42))
    add_text(slide, Inches(9.5), y + Inches(0.22), Inches(2.8), Inches(0.5), "📊 " + value, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 4: 产品概述 - 用户画像
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "01  产品概述 — 目标用户与核心价值")

users = [
    ("🧑‍💻 独立开发者", "统一展示 GitHub 项目、技术博客、技能树\n用 AI 分身处理客户咨询", C_PRIMARY),
    ("📱 自媒体创作者", "内容库统一管理草稿和发布\n查看跨平台数据分析", C_ACCENT),
    ("💕 恋爱中的情侣", "纪念日、旅行照片、心愿清单\n记录在专属空间", C_ROSE),
    ("📄 求职者", "高颜值在线简历 + 作品集\nAI 助手优化简历内容", C_EMERALD),
]
for i, (title, desc, color) in enumerate(users):
    x = Inches(0.6 + i * 3.1)
    card = add_rect(slide, x, Inches(1.3), Inches(2.9), Inches(2.5), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, Inches(1.3), Inches(2.9), Inches(0.06), color)
    add_text(slide, x + Inches(0.15), Inches(1.5), Inches(2.6), Inches(0.4), title, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_bullet_text(slide, x + Inches(0.2), Inches(2.1), Inches(2.5), Inches(1.5), desc.split("\n"), size=12, color=C_TEXT_SEC)

# 核心价值
add_text(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.4), "🎯 五大核心价值主张", size=20, color=C_WHITE, bold=True)
values = [
    ("🪪 个人 IP 一站式展示", "简历+技能树+作品集+社交矩阵"),
    ("📚 知识资产沉淀管理", "知识库+富文本+RAG 检索"),
    ("🤖 AI 分身能力赋能", "声音克隆+形象复刻+智能对话"),
    ("💝 生活情感数字珍藏", "恋爱记录+旅行地图+时光相册"),
    ("⚡ 工作效率加速器", "新媒体+电商+百宝箱集成"),
]
for i, (title, desc) in enumerate(values):
    x = Inches(0.6 + i * 2.5)
    add_rect(slide, x, Inches(4.8), Inches(2.3), Inches(1.8), RGBColor(0x25, 0x1D, 0x42))
    add_text(slide, x + Inches(0.1), Inches(4.9), Inches(2.1), Inches(0.8), title, size=13, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(5.6), Inches(2.1), Inches(0.8), desc, size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 5: 功能模块总览
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "02  核心功能模块 — 9 大模块全景图")

modules = [
    ("🏠", "个人主页", "Hero 区域 · 技能展示\n作品集 · 在线简历", C_BLUE),
    ("📚", "学习空间", "知识库 · 技能树\n富文本编辑 · AI 助手", C_EMERALD),
    ("💼", "工作助手", "新媒体运营 · 电商管理\n百宝箱工具", C_AMBER),
    ("🌸", "生活记录", "朋友圈 · 旅拍\n恋爱 · 音乐 · 电影", C_ROSE),
    ("🤖", "AI 分身", "声音克隆 · 形象复刻\n3D 角色 · RAG 对话", C_ACCENT),
    ("👤", "个人中心", "社交矩阵 · 留言墙\n数字名片 · VIP", C_PRIMARY),
    ("💰", "资产管理", "分类汇总 · 快速导航\n响应式网格", C_AMBER),
    ("🎨", "主题系统", "5 套主题一键切换\nCSS 变量驱动", C_PINK),
    ("🔐", "认证系统", "邮箱/用户名登录\nRLS 行级安全", C_EMERALD),
]
for i, (icon, title, desc, color) in enumerate(modules):
    row, col = divmod(i, 3)
    x = Inches(0.8 + col * 4.1)
    y = Inches(1.3 + row * 2.0)
    add_rect(slide, x, y, Inches(3.8), Inches(1.7), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, y, Inches(0.06), Inches(1.7), color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(0.5), Inches(0.5), icon, size=28)
    add_text(slide, x + Inches(0.8), y + Inches(0.12), Inches(2.5), Inches(0.35), title, size=16, color=C_WHITE, bold=True)
    add_bullet_text(slide, x + Inches(0.8), y + Inches(0.55), Inches(2.7), Inches(1.0), desc.split("\n"), size=11, color=C_TEXT_DIM)

# ══════════════════════════════════════════════
# SLIDE 6: 个人主页模块
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "02  个人主页 — 一站式个人 IP 展示")
add_text(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.4), "🏠 个人主页模块", size=22, color=C_WHITE, bold=True)
features_home = [
    "✦ Hero 区域：头像 + 个人简介 + 在线简历切换",
    "✦ 技能展示：分类技能条形图（ECharts 可视化）",
    "✦ 项目作品集：卡片网格布局，封面图 + 标签",
    "✦ 兴趣爱好：3D 卡片轮播，悬停展开手势滑动",
    "✦ 自媒体矩阵：各平台社交链接聚合展示",
    "✦ 联系区域：邮件 / 社交媒体 / 二维码",
    "✦ 编辑模式：实时编辑 + 撤销 / 重做 + PDF 导出",
]
add_bullet_text(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(4), features_home, size=13, color=C_TEXT_SEC)
# 截图
add_screenshot(slide, "app_homepage_dashboard_1777195791085.png", Inches(6.2), Inches(1.1), Inches(6.5), Inches(3.7))
# 流程图
add_rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5), RGBColor(0x1A, 0x14, 0x30))
add_text(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.3), "交互流程", size=14, color=C_PRIMARY, bold=True)
add_text(slide, Inches(0.8), Inches(5.95), Inches(11.5), Inches(0.9),
         "页面加载 → 读取 ThemeContext → 渲染 Hero/Skills/Portfolio → 点击「查看简历」→ 切换 OnlineResume 视图 → 点击主题按钮 → 全局色调切换",
         size=12, color=C_TEXT_DIM)

# ══════════════════════════════════════════════
# SLIDE 7: 学习空间模块
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "02  学习空间 — 知识沉淀与智能问答")
# 截图
add_screenshot(slide, "learning_workspace_1777195794905.png", Inches(0.6), Inches(1.1), Inches(7), Inches(4))
# 右侧功能
add_text(slide, Inches(8), Inches(1.1), Inches(5), Inches(0.4), "📚 学习空间模块", size=22, color=C_WHITE, bold=True)
features_learn = [
    "✦ 知识库视图：三栏布局（分类/列表/详情）",
    "✦ 技能树视图：可视化技能成长路径",
    "✦ 富文本编辑器：基于 Tiptap，支持图片/链接",
    "✦ 文档导入：本地文档批量解析导入",
    "✦ AI 阅读助手：RAG 引擎智能问答浮窗",
    "✦ 在线简历：结构化简历编辑 + PDF 导出",
    "✦ 文章 CRUD：新建/编辑/删除，标签分类",
]
add_bullet_text(slide, Inches(8), Inches(1.6), Inches(4.8), Inches(3.5), features_learn, size=13, color=C_TEXT_SEC)
# RAG 引擎
add_rect(slide, Inches(8), Inches(5.2), Inches(4.8), Inches(2), RGBColor(0x25, 0x1D, 0x42))
add_text(slide, Inches(8.2), Inches(5.3), Inches(4), Inches(0.3), "🧠 自研 RAG 知识检索引擎", size=14, color=C_ACCENT, bold=True)
add_bullet_text(slide, Inches(8.2), Inches(5.7), Inches(4.4), Inches(1.3), [
    "关键词匹配 + 语义理解双引擎",
    "覆盖学习/生活/娱乐全模块知识",
    "多种回复模式（推荐/方法/介绍）",
    "本地运行，无需外部 API 调用",
], size=11, color=C_TEXT_DIM)

# ══════════════════════════════════════════════
# SLIDE 8: 工作助手模块
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "02  工作助手 — 新媒体 · 电商 · 百宝箱")
add_text(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.5), "💼 工作助手模块", size=22, color=C_WHITE, bold=True)
# 三个子模块
subs = [
    ("📱 新媒体运营", ["内容库管理（文章/视频）", "发布管理（草稿/定时/已发）", "数据分析（四维统计看板）", "平台筛选（微信/微博/小红书/抖音）"], C_PRIMARY),
    ("🛒 电商运营", ["数据概览（销售额/订单/转化率）", "商品管理（上下架/定价/库存）", "订单管理（全流程追踪）", "经营分析（可视化图表）"], C_AMBER),
    ("🧰 百宝箱", ["工具收藏管理", "分类搜索过滤", "快速导航入口", "开发/设计/学习资源"], C_EMERALD),
]
for i, (title, items, color) in enumerate(subs):
    x = Inches(0.6 + i * 4.15)
    add_rect(slide, x, Inches(1.7), Inches(3.9), Inches(2.8), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, Inches(1.7), Inches(3.9), Inches(0.06), color)
    add_text(slide, x + Inches(0.15), Inches(1.85), Inches(3.5), Inches(0.35), title, size=16, color=color, bold=True)
    add_bullet_text(slide, x + Inches(0.15), Inches(2.3), Inches(3.5), Inches(2.0), items, size=12, color=C_TEXT_SEC)
# 截图
add_screenshot(slide, "work_assistant_page_1777195809723.png", Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.5))

# ══════════════════════════════════════════════
# SLIDE 9: 生活记录模块
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "02  生活记录 — 记录生活中的美好时刻")
# 截图
add_screenshot(slide, "life_record_page_1777195819695.png", Inches(0.6), Inches(1.1), Inches(6.5), Inches(3.7))
# 右侧
add_text(slide, Inches(7.5), Inches(1.1), Inches(5), Inches(0.4), "🌸 生活记录模块", size=22, color=C_WHITE, bold=True)
life_modules = [
    ("💕 恋爱记录", "倒计时 · 时间轴事件 · 时光相册 · 许愿清单 · 祝福墙"),
    ("📸 旅拍相册", "中国地图轨迹可视化 · 旅行照片管理 · 详情弹窗"),
    ("💬 朋友圈", "图文动态 · 点赞评论 · 语音输入"),
    ("🎵 音乐墙", "音乐收藏展示"),
    ("🎬 收藏电影", "豆瓣风格电影收藏墙"),
    ("💪 运动", "运动数据记录"),
    ("🎮 游戏", "游戏收藏管理"),
]
for i, (title, desc) in enumerate(life_modules):
    y = Inches(1.6 + i * 0.72)
    add_text(slide, Inches(7.5), y, Inches(2), Inches(0.3), title, size=13, color=C_ROSE, bold=True)
    add_text(slide, Inches(9.2), y, Inches(3.5), Inches(0.6), desc, size=11, color=C_TEXT_DIM)

# ══════════════════════════════════════════════
# SLIDE 10: AI 分身模块
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "02  AI 分身 — 三步创建专属数字分身")
# 截图
add_screenshot(slide, "ai_character_page_1777195834660.png", Inches(0.6), Inches(1.1), Inches(7), Inches(4))
# 右侧步骤
add_text(slide, Inches(8), Inches(1.1), Inches(5), Inches(0.4), "🤖 AI 分身模块", size=22, color=C_WHITE, bold=True)
steps = [
    ("Step 1", "🎙️ 声音克隆", "录制/上传音频样本\n生成个性化声音模型", C_PRIMARY),
    ("Step 2", "👤 形象复刻", "选择/上传头像\n配置卡通/写实风格", C_ACCENT),
    ("Step 3", "💬 AI 对话", "Three.js 3D 角色渲染\n多模态输入 + RAG 智能回复", C_EMERALD),
]
for i, (step, title, desc, color) in enumerate(steps):
    y = Inches(1.7 + i * 1.5)
    add_rect(slide, Inches(8), y, Inches(4.8), Inches(1.3), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, Inches(8), y, Inches(0.06), Inches(1.3), color)
    add_text(slide, Inches(8.2), y + Inches(0.05), Inches(1), Inches(0.3), step, size=11, color=color, bold=True)
    add_text(slide, Inches(8.2), y + Inches(0.35), Inches(4.3), Inches(0.3), title, size=16, color=C_WHITE, bold=True)
    add_text(slide, Inches(8.2), y + Inches(0.7), Inches(4.3), Inches(0.5), desc, size=11, color=C_TEXT_DIM)

# 关键特色
add_rect(slide, Inches(0.6), Inches(5.5), Inches(12), Inches(1.5), RGBColor(0x1A, 0x14, 0x30))
add_text(slide, Inches(0.8), Inches(5.6), Inches(3), Inches(0.3), "关键技术特色", size=14, color=C_ACCENT, bold=True)
techs = ["Three.js 3D 角色渲染", "Web Speech API 语音朗读", "RAG 引擎知识检索", "多模态输入（文字/图片/语音）"]
for i, t in enumerate(techs):
    x = Inches(0.8 + i * 3)
    add_rect(slide, x, Inches(6.0), Inches(2.7), Inches(0.45), RGBColor(0x25, 0x1D, 0x42))
    add_text(slide, x, Inches(6.03), Inches(2.7), Inches(0.4), t, size=11, color=C_PRIMARY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 11: 技术架构亮点
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "03  技术架构亮点")

# 前端架构
add_text(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.4), "📱 前端技术栈", size=18, color=C_PRIMARY, bold=True)
fe_techs = [
    ("React 19.2.0", "声明式 UI · Hooks · Concurrent", C_BLUE),
    ("TypeScript 5.9", "全覆盖类型安全 · 接口先行", C_BLUE),
    ("Vite 7.3.1", "极速构建 · HMR · 模块联邦", C_AMBER),
    ("Tailwind CSS 4.x", "原子化 CSS · JIT 编译", C_PRIMARY),
    ("Framer Motion", "页面过渡 · 微交互动效", C_ACCENT),
    ("Three.js", "3D 角色渲染 · WebGL", C_EMERALD),
]
for i, (name, desc, color) in enumerate(fe_techs):
    row, col = divmod(i, 2)
    x = Inches(0.6 + col * 3)
    y = Inches(1.6 + row * 0.75)
    add_rect(slide, x, y, Inches(2.8), Inches(0.6), RGBColor(0x1A, 0x14, 0x30))
    add_text(slide, x + Inches(0.1), y + Inches(0.03), Inches(2.6), Inches(0.3), name, size=12, color=color, bold=True)
    add_text(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.6), Inches(0.25), desc, size=9, color=C_TEXT_DIM)

# 后端架构
add_text(slide, Inches(7), Inches(1.1), Inches(5), Inches(0.4), "💎 后端 & AI", size=18, color=C_ACCENT, bold=True)
be_techs = [
    ("Supabase BaaS", "云端 PostgreSQL + Auth + RLS", C_EMERALD),
    ("Tiptap 3.20", "富文本编辑器 · 知识写作", C_AMBER),
    ("ECharts 6.0", "数据可视化 · 中国地图", C_BLUE),
    ("RAG Engine", "自研知识检索 · 多模式回复", C_ACCENT),
    ("html2canvas", "数字名片截图生成", C_ROSE),
    ("Web Speech API", "语音合成 · 声音朗读", C_PRIMARY),
]
for i, (name, desc, color) in enumerate(be_techs):
    row, col = divmod(i, 2)
    x = Inches(6.8 + col * 3)
    y = Inches(1.6 + row * 0.75)
    add_rect(slide, x, y, Inches(2.8), Inches(0.6), RGBColor(0x1A, 0x14, 0x30))
    add_text(slide, x + Inches(0.1), y + Inches(0.03), Inches(2.6), Inches(0.3), name, size=12, color=color, bold=True)
    add_text(slide, x + Inches(0.1), y + Inches(0.3), Inches(2.6), Inches(0.25), desc, size=9, color=C_TEXT_DIM)

# 架构设计原则
add_text(slide, Inches(0.8), Inches(4.1), Inches(11), Inches(0.4), "📐 架构设计原则", size=18, color=C_WHITE, bold=True)
principles = [
    ("模块化", "9 大页面 + 30+ 组件\n按功能域拆分"),
    ("类型安全", "TypeScript 全覆盖\n接口先行定义"),
    ("状态管理", "Context 轻量级\nTheme/Auth/Settings"),
    ("主题驱动", "CSS 变量 + Config\n40+ 配色变量"),
    ("渐进增强", "Supabase 失败\n自动回退 Mock"),
    ("响应式", "移动端底部导航\n桌面端顶部导航"),
]
for i, (title, desc) in enumerate(principles):
    x = Inches(0.6 + i * 2.1)
    add_rect(slide, x, Inches(4.6), Inches(1.9), Inches(2.3), RGBColor(0x1A, 0x14, 0x30))
    colors_p = [C_PRIMARY, C_ACCENT, C_EMERALD, C_AMBER, C_ROSE, C_BLUE]
    add_rect(slide, x, Inches(4.6), Inches(1.9), Inches(0.06), colors_p[i])
    add_text(slide, x + Inches(0.1), Inches(4.75), Inches(1.7), Inches(0.35), title, size=14, color=colors_p[i], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(5.2), Inches(1.7), Inches(1.2), desc, size=11, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 12: 用户体验设计
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "04  用户体验设计 — 5 套主题系统")

themes = [
    ("☀️ 极简", "白色背景 · 蓝色主色\n简洁线条 · 专业求职", RGBColor(0x25, 0x63, 0xEB)),
    ("🌙 赛博", "暗黑霓虹 · 青蓝紫色\n赛博朋克 · 科技极客", RGBColor(0x00, 0xD4, 0xFF)),
    ("🎨 艺术", "暖橙珊瑚 · 手工艺感\n柔和渐变 · 创意设计", RGBColor(0xE8, 0x5D, 0x75)),
    ("🌈 童趣", "粉色明黄 · 圆润边角\n明亮活泼 · 生活博主", RGBColor(0xF4, 0x72, 0xB6)),
    ("📜 复古", "深金棕红 · 老式排版\n暖黄棕色 · 文学怀旧", RGBColor(0x8B, 0x69, 0x14)),
]
for i, (icon_name, desc, color) in enumerate(themes):
    x = Inches(0.5 + i * 2.5)
    add_rect(slide, x, Inches(1.3), Inches(2.3), Inches(2.5), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, Inches(1.3), Inches(2.3), Inches(0.08), color)
    add_text(slide, x, Inches(1.5), Inches(2.3), Inches(0.5), icon_name, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.15), Inches(2.1), Inches(2.0), Inches(1.5), desc, size=11, color=C_TEXT_SEC, align=PP_ALIGN.CENTER)

# 主题截图对比
add_screenshot(slide, "social_profile_logged_in_1777195873008.png", Inches(0.6), Inches(4.1), Inches(5.8), Inches(3.2))
add_screenshot(slide, "cyberpunk_theme_profile_1777195906398.png", Inches(6.8), Inches(4.1), Inches(5.8), Inches(3.2))
add_text(slide, Inches(0.6), Inches(3.8), Inches(5.8), Inches(0.3), "极简主题", size=12, color=C_BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(6.8), Inches(3.8), Inches(5.8), Inches(0.3), "赛博主题", size=12, color=C_PRIMARY, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 13: 市场分析
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "05  市场分析 — 竞品对比与商业机会")

# 竞品对比表
add_text(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.4), "🔍 竞品分析矩阵", size=18, color=C_WHITE, bold=True)
competitors = [
    ("Notion", "知识管理", "✅ 生态完善", "❌ 无AI分身/个人主页/生活记录"),
    ("Linktree", "链接聚合", "✅ 简单易用", "❌ 功能单一/无知识库/无主题"),
    ("Character.AI", "AI角色", "✅ 角色丰富", "❌ 无个人知识库/无资产管理"),
    ("飞书文档", "协作文档", "✅ 企业级协作", "❌ 偏B端/无个人IP展示"),
    ("GitHub Pages", "开发者主页", "✅ 开发者生态", "❌ 技术门槛高/功能单一"),
]
# Header
add_rect(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.4), C_ACCENT)
headers = ["竞品", "定位", "优势", "不足"]
hx = [0.7, 2.2, 4, 7]
hw = [1.3, 1.5, 2.8, 5]
for j, h in enumerate(headers):
    add_text(slide, Inches(hx[j]), Inches(1.52), Inches(hw[j]), Inches(0.35), h, size=12, color=C_WHITE, bold=True)
for i, (name, pos, adv, dis) in enumerate(competitors):
    y = Inches(1.95 + i * 0.55)
    bg_c = RGBColor(0x1A, 0x14, 0x30) if i % 2 == 0 else RGBColor(0x22, 0x1A, 0x3D)
    add_rect(slide, Inches(0.6), y, Inches(12), Inches(0.5), bg_c)
    vals = [name, pos, adv, dis]
    for j, v in enumerate(vals):
        add_text(slide, Inches(hx[j]), y + Inches(0.05), Inches(hw[j]), Inches(0.4), v, size=11, color=C_TEXT_SEC)

# Opclaw 差异化
add_rect(slide, Inches(0.6), Inches(4.75), Inches(12), Inches(0.5), C_PRIMARY)
add_text(slide, Inches(0.7), Inches(4.77), Inches(11.8), Inches(0.45),
         "🌈 Opclaw：市场唯一「个人展示 + 知识管理 + AI 分身 + 生活记录 + 工作助手」五合一平台",
         size=13, color=C_BG_DARK, bold=True)

# 商业模式
add_text(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.4), "💰 商业模式", size=18, color=C_WHITE, bold=True)
plans = [
    ("🆓 免费版", "¥0", "基础 5 模块 + 5 套主题\n1 个 AI 分身 + 1GB 存储", C_EMERALD),
    ("💎 VIP 版", "¥19.9/月", "全部功能 + 高级主题\n5 个 AI 分身 + 10GB", C_AMBER),
    ("👑 Pro 版", "¥49.9/月", "无限 AI 分身 + 存储\nAPI 接入 + 白标定制", C_ACCENT),
]
for i, (name, price, desc, color) in enumerate(plans):
    x = Inches(0.6 + i * 4.15)
    add_rect(slide, x, Inches(5.9), Inches(3.9), Inches(1.3), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, Inches(5.9), Inches(3.9), Inches(0.06), color)
    add_text(slide, x + Inches(0.15), Inches(6.0), Inches(1.5), Inches(0.3), name, size=14, color=color, bold=True)
    add_text(slide, x + Inches(2.2), Inches(6.0), Inches(1.5), Inches(0.3), price, size=16, color=C_WHITE, bold=True, align=PP_ALIGN.RIGHT)
    add_text(slide, x + Inches(0.15), Inches(6.35), Inches(3.5), Inches(0.8), desc, size=11, color=C_TEXT_DIM)

# ══════════════════════════════════════════════
# SLIDE 14: 评分亮点
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "05  产品评分 — 多维度综合评估")

scores = [
    ("视觉设计", "9.5", "5套主题·玻璃拟态·流畅动画", C_ROSE),
    ("AI 创新", "9.0", "RAG引擎·声音克隆·3D角色", C_ACCENT),
    ("架构质量", "8.5", "React19·TS全覆盖·模块化", C_BLUE),
    ("功能完整", "8.5", "9大模块覆盖全场景", C_EMERALD),
    ("用户体验", "8.5", "响应式·动画·星光特效", C_AMBER),
    ("可扩展性", "9.0", "主题/模块/RAG可扩展", C_PRIMARY),
    ("技术前瞻", "9.0", "React19·Vite7·Three.js", C_PINK),
    ("差异化", "9.5", "五合一平台·市场独一无二", C_ROSE),
]
for i, (dim, score, desc, color) in enumerate(scores):
    row, col = divmod(i, 4)
    x = Inches(0.6 + col * 3.1)
    y = Inches(1.3 + row * 2.8)
    add_rect(slide, x, y, Inches(2.9), Inches(2.5), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, y, Inches(2.9), Inches(0.06), color)
    add_text(slide, x, y + Inches(0.2), Inches(2.9), Inches(0.4), dim, size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.7), Inches(2.9), Inches(0.8), score, size=40, color=color, bold=True, align=PP_ALIGN.CENTER)
    # 分数条
    bar_w = float(score) / 10 * 2.5
    add_rect(slide, x + Inches(0.2), y + Inches(1.6), Inches(2.5), Inches(0.12), RGBColor(0x25, 0x1D, 0x42))
    add_gradient_rect(slide, x + Inches(0.2), y + Inches(1.6), Inches(bar_w), Inches(0.12), color, RGBColor(0x25, 0x1D, 0x42))
    add_text(slide, x + Inches(0.1), y + Inches(1.85), Inches(2.7), Inches(0.5), desc, size=10, color=C_TEXT_DIM, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════
# SLIDE 15: 总结与展望
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
section_title_bar(slide, "06  总结与展望")

# 核心优势
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4), "🏆 核心优势", size=20, color=C_WHITE, bold=True)
advantages = [
    ("🎯 定位独特", "市场唯一五合一平台", C_PRIMARY),
    ("🎨 视觉出众", "5套主题·CSS变量驱动", C_ACCENT),
    ("🤖 AI 创新", "RAG+声音克隆+3D形象", C_EMERALD),
    ("🏗️ 架构扎实", "React19+TS+模块化", C_AMBER),
    ("✨ 体验流畅", "动画+响应式+特效", C_ROSE),
    ("🔧 可扩展", "模块/主题/RAG扩展", C_BLUE),
]
for i, (title, desc, color) in enumerate(advantages):
    row, col = divmod(i, 3)
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.6 + row * 1.2)
    add_rect(slide, x, y, Inches(3.8), Inches(1.0), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, y, Inches(0.06), Inches(1.0), color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.35), title, size=14, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(3.4), Inches(0.4), desc, size=12, color=C_TEXT_DIM)

# 版本路线图
add_text(slide, Inches(0.8), Inches(4.2), Inches(10), Inches(0.4), "🗺️ 版本路线图", size=18, color=C_WHITE, bold=True)
roadmap = [
    ("v1.0-v1.5", "✅ 已完成", "基础框架·个人主页·学习空间·生活记录\nAI分身·个人中心·Supabase认证", C_EMERALD),
    ("v2.0", "🚧 规划中", "AI简历优化器·实时协作编辑\n数据可视化大屏·语音交互系统", C_AMBER),
    ("v3.0", "🔬 研究中", "PWA离线支持·WebGL 3D空间\n企业版·API开放平台", C_ACCENT),
]
for i, (ver, status, desc, color) in enumerate(roadmap):
    x = Inches(0.6 + i * 4.15)
    add_rect(slide, x, Inches(4.7), Inches(3.9), Inches(2.3), RGBColor(0x1A, 0x14, 0x30))
    add_rect(slide, x, Inches(4.7), Inches(3.9), Inches(0.06), color)
    add_text(slide, x + Inches(0.15), Inches(4.85), Inches(1.5), Inches(0.3), ver, size=18, color=color, bold=True)
    add_text(slide, x + Inches(2.2), Inches(4.85), Inches(1.5), Inches(0.3), status, size=12, color=C_TEXT_SEC, align=PP_ALIGN.RIGHT)
    add_text(slide, x + Inches(0.15), Inches(5.3), Inches(3.5), Inches(1.5), desc, size=12, color=C_TEXT_DIM)

# ══════════════════════════════════════════════
# SLIDE 16: 结尾页
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_gradient_rect(slide, Inches(0), Inches(0), W, Inches(0.06), C_PRIMARY, C_ACCENT)
add_gradient_rect(slide, Inches(0), Inches(7.44), W, Inches(0.06), C_ACCENT, C_PRIMARY)
add_text(slide, Inches(0), Inches(1.5), W, Inches(0.8), "🌈", size=60, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(2.5), W, Inches(0.8), "Thank You", size=48, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0), Inches(3.3), W, Inches(0.6), "一个人的数字宇宙，从这里开始！", size=22, color=C_PRIMARY, align=PP_ALIGN.CENTER)
# 联系信息
add_rect(slide, Inches(4), Inches(4.3), Inches(5.3), Inches(2), RGBColor(0x1A, 0x14, 0x30))
info_items = [
    "👤  晓叶  |  全栈开发工程师 & AI 研究员",
    "📍  中国 · 杭州",
    "📧  wyxcode@qq.com",
    "🌐  https://xiaoyu.dev",
]
add_bullet_text(slide, Inches(4.5), Inches(4.5), Inches(4.5), Inches(1.8), info_items, size=13, color=C_TEXT_SEC)

# ── 保存 ──
prs.save(OUTPUT)
print(f"\n✅ PPT 已成功生成: {OUTPUT}")
print(f"   共 {len(prs.slides)} 页幻灯片")
